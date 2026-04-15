"""
Generates two files inside the Presentation/ folder:
  1. Presentation_Script.docx  – full 5-person speaking script
  2. Presentation_Slides.pptx  – slide deck (key points per slide)
"""

# ── WORD DOCUMENT ────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import copy

doc = Document()

# ── Page margins ──────────────────────────────────────────────────────────────
section = doc.sections[0]
section.top_margin    = Cm(2.0)
section.bottom_margin = Cm(2.0)
section.left_margin   = Cm(2.5)
section.right_margin  = Cm(2.5)

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x3C, 0x6D)   # system primary blue
DARK   = RGBColor(0x0F, 0x17, 0x2A)   # near-black
MUTED  = RGBColor(0x64, 0x74, 0x8B)   # slate-500
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PERSON_COLORS = [
    RGBColor(0x1A, 0x3C, 0x6D),  # P1 – navy
    RGBColor(0x16, 0x5F, 0x9F),  # P2 – blue
    RGBColor(0x0E, 0x7C, 0x65),  # P3 – teal
    RGBColor(0x7C, 0x2D, 0x12),  # P4 – brown
    RGBColor(0x4C, 0x1D, 0x95),  # P5 – purple
]

# ── Helper functions ──────────────────────────────────────────────────────────
def add_paragraph(text="", bold=False, italic=False, size=11,
                  color=None, align=WD_ALIGN_PARAGRAPH.LEFT, space_before=0, space_after=6):
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after  = Pt(space_after)
    if text:
        run = p.add_run(text)
        run.bold   = bold
        run.italic = italic
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color
    return p

def add_heading(text, level=1, color=NAVY):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(14 if level == 1 else 10)
    p.paragraph_format.space_after  = Pt(4)
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt({1: 20, 2: 15, 3: 13, 4: 11}[level])
    run.font.color.rgb = color
    return p

def add_rule():
    """Horizontal line via bottom border on a paragraph."""
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after  = Pt(4)
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '6')
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), '1A3C6D')
    pBdr.append(bottom)
    pPr.append(pBdr)
    return p

def add_bullet(text, indent=0):
    p = doc.add_paragraph(style='List Bullet')
    p.paragraph_format.left_indent   = Inches(0.3 + indent * 0.25)
    p.paragraph_format.space_after   = Pt(3)
    run = p.add_run(text)
    run.font.size = Pt(10.5)
    return p

def add_script_block(text, person_idx=0):
    """Quoted speaking block with left border colour."""
    for line in text.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        p = doc.add_paragraph()
        p.paragraph_format.left_indent  = Inches(0.35)
        p.paragraph_format.space_after  = Pt(5)
        run = p.add_run(line)
        run.font.size = Pt(10.5)
        run.italic = True
        # left border
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        left = OxmlElement('w:left')
        left.set(qn('w:val'), 'single')
        left.set(qn('w:sz'), '12')
        left.set(qn('w:space'), '8')
        clr = PERSON_COLORS[person_idx]
        left.set(qn('w:color'), f'{clr[0]:02X}{clr[1]:02X}{clr[2]:02X}')
        pBdr.append(left)
        pPr.append(pBdr)

def person_banner(number, name, role, color):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after  = Pt(6)
    shading = OxmlElement('w:shd')
    clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
    shading.set(qn('w:val'),   'clear')
    shading.set(qn('w:color'), 'auto')
    shading.set(qn('w:fill'),  clr_hex)
    pPr = p._p.get_or_add_pPr()
    pPr.append(shading)
    run = p.add_run(f'  PERSON {number}  ·  {name}  ·  {role}')
    run.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE


# ════════════════════════════════════════════════════════════════════════════
#  COVER PAGE
# ════════════════════════════════════════════════════════════════════════════
add_paragraph()
add_paragraph()
add_paragraph(
    'UNIVERSITY SEAT RESERVATION SYSTEM',
    bold=True, size=24, color=NAVY,
    align=WD_ALIGN_PARAGRAPH.CENTER, space_before=40
)
add_paragraph(
    'Group Presentation — Full Speaking Script',
    size=14, color=MUTED,
    align=WD_ALIGN_PARAGRAPH.CENTER, space_before=6
)
add_rule()
add_paragraph()

# Team table
table = doc.add_table(rows=6, cols=3)
table.style = 'Table Grid'
hdr = table.rows[0].cells
for cell, txt in zip(hdr, ['Person', 'Role', 'Section']):
    cell.paragraphs[0].clear()
    run = cell.paragraphs[0].add_run(txt)
    run.bold = True
    run.font.color.rgb = WHITE
    cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    shading = OxmlElement('w:shd')
    shading.set(qn('w:val'),   'clear')
    shading.set(qn('w:color'), 'auto')
    shading.set(qn('w:fill'),  '1A3C6D')
    cell._tc.get_or_add_tcPr().append(shading)

rows_data = [
    ('Person 1', 'Project Lead',          'Introduction & System Overview'),
    ('Person 2', 'Frontend Developer',    'Pages, Components & UI'),
    ('Person 3', 'Backend Developer',     'PHP APIs & Server Logic'),
    ('Person 4', 'Database Admin',        'MySQL Schema & Data Flow'),
    ('Person 5', 'Systems Integrator',    'Full Flow & Conclusion'),
]
for i, (p, r, s) in enumerate(rows_data):
    row = table.rows[i + 1]
    row.cells[0].text = p
    row.cells[1].text = r
    row.cells[2].text = s
    clr_hex = f'{PERSON_COLORS[i][0]:02X}{PERSON_COLORS[i][1]:02X}{PERSON_COLORS[i][2]:02X}'
    for cell in row.cells:
        shading = OxmlElement('w:shd')
        shading.set(qn('w:val'),   'clear')
        shading.set(qn('w:color'), 'auto')
        shading.set(qn('w:fill'),  'F0F4FF' if i % 2 == 0 else 'FFFFFF')
        cell._tc.get_or_add_tcPr().append(shading)

add_paragraph()
add_paragraph(
    'Demo Login:  STU001  /  password123      |      Start:  npm run dev',
    size=10, color=MUTED, align=WD_ALIGN_PARAGRAPH.CENTER
)
doc.add_page_break()

# ════════════════════════════════════════════════════════════════════════════
#  QUICK FACTS BOX
# ════════════════════════════════════════════════════════════════════════════
add_heading('Quick Reference Facts', 2)
facts = [
    ('Frontend',        'React 19 + Vite 8 + Tailwind CSS'),
    ('Backend',         'PHP 8.0 (PHP built-in server via XAMPP)'),
    ('Database',        'MySQL — managed via phpMyAdmin'),
    ('Authentication',  'PHP sessions + bcrypt password hashing'),
    ('Routing',         '4 routes: / · /login · /signup · /home'),
    ('Venue zones',     'Shade A (6 tables) · Shade B (25 tables)'),
    ('Capacity',        '6 people per table'),
    ('User roles',      'Student · Staff · Guest'),
    ('Password rules',  '8+ characters · 1 uppercase · 1 special character'),
    ('Start command',   'npm run dev  (from project root)'),
    ('Frontend URL',    'http://localhost:5173'),
    ('Backend URL',     'http://localhost:8000'),
]
t = doc.add_table(rows=len(facts), cols=2)
t.style = 'Table Grid'
for i, (k, v) in enumerate(facts):
    t.rows[i].cells[0].text = k
    t.rows[i].cells[1].text = v
    run0 = t.rows[i].cells[0].paragraphs[0].runs[0]
    run0.bold = True
    run0.font.color.rgb = NAVY
    bg = 'EFF6FF' if i % 2 == 0 else 'FFFFFF'
    for cell in t.rows[i].cells:
        shading = OxmlElement('w:shd')
        shading.set(qn('w:val'),   'clear')
        shading.set(qn('w:color'), 'auto')
        shading.set(qn('w:fill'),  bg)
        cell._tc.get_or_add_tcPr().append(shading)
add_paragraph()
doc.add_page_break()

# ════════════════════════════════════════════════════════════════════════════
#  PERSON 1 — INTRODUCTION
# ════════════════════════════════════════════════════════════════════════════
person_banner(1, 'Project Lead', 'Introduction & System Overview', PERSON_COLORS[0])

add_heading('What to Show on Screen', 3, PERSON_COLORS[0])
add_bullet('Running app at http://localhost:5173 (landing page visible)')
add_bullet('Have the dashboard (/home) ready to switch to')

add_heading('Speaking Script', 3, PERSON_COLORS[0])
add_script_block("""
Good [morning/afternoon], everyone. My name is [Name], and on behalf of our team, welcome to our project presentation.

The system we built is called the University Seat Reservation System. The idea is simple — if you've ever walked onto campus looking for a table to sit and study, only to find everything is taken, this system solves that problem.

Instead of wandering around, students, staff, and guests can open this web app, see which tables are available right now, pick a spot, choose a time, and confirm their booking — all in under a minute.
""", 0)

add_paragraph('[Demo the live app while talking]', italic=True, size=10, color=MUTED)

add_script_block("""
This is the homepage — the first thing you see when you visit. It explains what the system does and gives you a button to sign in or create an account.

When you log in, you land on the dashboard — this is where the real action happens. You can see a map of all the tables on campus, split into two zones: Shade A and Shade B. Each table shows how many seats are still available. You click a table, fill in your group size and time, and that's your reservation confirmed.

The system supports three types of users — students, staff, and guests — and each has slightly different access rules. For example, guests can only book in Shade B, and they're limited to a group of three.

The frontend is built with React and Vite. The backend is PHP, running through XAMPP. The database is MySQL. Everything is connected: when you click a button in the browser, JavaScript sends a request to PHP, PHP talks to MySQL, and the result comes back to your screen.

I'll now hand it over to [Person 2's name], who will take you through the frontend in detail.
""", 0)

add_rule()

# ════════════════════════════════════════════════════════════════════════════
#  PERSON 2 — FRONTEND
# ════════════════════════════════════════════════════════════════════════════
doc.add_page_break()
person_banner(2, 'Frontend Developer', 'Pages, Components & UI', PERSON_COLORS[1])

add_heading('Files to Navigate (in order)', 3, PERSON_COLORS[1])
add_bullet('frontend/src/main.jsx')
add_bullet('frontend/src/Login.jsx')
add_bullet('frontend/src/Signup.jsx')
add_bullet('frontend/src/Home.jsx')
add_bullet('frontend/src/Components/SeatMap.jsx')

add_heading('Speaking Script', 3, PERSON_COLORS[1])

add_paragraph('Open: frontend/src/main.jsx', bold=True, size=10, color=PERSON_COLORS[1])
add_script_block("""
Thank you, [Person 1]. My name is [Name], and I'm going to walk you through the frontend — everything the user sees and interacts with.

Our frontend is a React application. React lets us build the interface as a collection of small, reusable pieces called components. Each component handles one specific part of the page.

Let's start with the entry point — main.jsx. This file sets up the routing — it decides which page to show depending on the URL. We have four routes: the slash for the landing page, /login, /signup, and /home for the main dashboard. When you type a URL, React looks at this file and shows the right component. It's all handled client-side — no full page reloads, everything feels instant.
""", 1)

add_paragraph('Open: frontend/src/Login.jsx', bold=True, size=10, color=PERSON_COLORS[1])
add_script_block("""
Now the login page. It's a form with Student ID and password. When submitted, it calls handleLogin, which uses our API layer to send credentials to the backend.

We added a live password checklist — it appears when the user focuses on the password field and checks three rules in real time: minimum 8 characters, one uppercase letter, one special character.

On the login page this checklist is informational only — it doesn't block submission. This is intentional: if someone registered before we introduced the new rules, we don't want to lock them out. The server decides whether credentials are valid.
""", 1)

add_paragraph('Open: frontend/src/Signup.jsx', bold=True, size=10, color=PERSON_COLORS[1])
add_script_block("""
On the signup page the rules are enforced. The submit button is actually disabled until all three password conditions are green. You can see the passwordValid check — if any rule hasn't passed and the user has touched the field, the button grays out and you cannot submit.

The role dropdown lets users identify as student, staff, or guest, and this affects what they can do inside the system.
""", 1)

add_paragraph('Open: frontend/src/Home.jsx', bold=True, size=10, color=PERSON_COLORS[1])
add_script_block("""
The most important page is Home.jsx — the dashboard. The first thing it does is call auth.me() — it asks the backend whether this user is still logged in. If not, it immediately redirects to the homepage.

Once confirmed, it loads two things simultaneously: the list of venues and the user's reservations — using Promise.all so both load in parallel.

The layout is a two-column grid. On the left is the seat map. On the right is the booking form and the reservations list.
""", 1)

add_paragraph('Open: frontend/src/Components/SeatMap.jsx', bold=True, size=10, color=PERSON_COLORS[1])
add_script_block("""
The SeatMap component shows all the tables visually. It receives venue data and splits it into Shade A and Shade B.

Each table is a button. The colour tells you its state: white means available, blue means selected by you, grey means full, and amber means restricted — shown to guests for Shade A tables.

There's also a small progress bar on each card showing how many of the six seats are already booked.

Our frontend also uses Tailwind CSS for styling — a utility-first framework that lets us apply styles directly in the JSX using class names.

I'll now hand over to [Person 3], who will explain what happens on the server side when those API calls go out.
""", 1)

add_rule()

# ════════════════════════════════════════════════════════════════════════════
#  PERSON 3 — BACKEND
# ════════════════════════════════════════════════════════════════════════════
doc.add_page_break()
person_banner(3, 'Backend Developer', 'PHP APIs & Server Logic', PERSON_COLORS[2])

add_heading('Files to Navigate (in order)', 3, PERSON_COLORS[2])
add_bullet('backend/config/database.php')
add_bullet('backend/api/auth/login.php')
add_bullet('backend/api/auth/register.php')
add_bullet('backend/api/reservations/create.php')

add_heading('Speaking Script', 3, PERSON_COLORS[2])

add_paragraph('Open: backend/config/database.php', bold=True, size=10, color=PERSON_COLORS[2])
add_script_block("""
Thank you, [Person 2]. My name is [Name] and I'll explain the backend — the server-side code that handles all the logic, security, and data processing.

Every single PHP file in our backend starts by including one file: database.php. It does several important things.

First, it loads our environment variables from a .env file. Instead of putting the database password directly in the code — which would be a security problem — we keep it in a separate file that is not committed to GitHub.

Second, it sets the CORS headers. CORS is a browser security rule. Our frontend runs on port 5173 and our backend on port 8000. Without these headers, the browser would block all our API requests.

Third, it provides four helper functions every other file uses: getDb() opens the database connection, jsonResponse() sends JSON back to the frontend, errorResponse() sends error messages, and requireAuth() checks if the user is logged in. This design means we don't repeat ourselves — every file just calls requireAuth() at the top.
""", 2)

add_paragraph('Open: backend/api/auth/login.php', bold=True, size=10, color=PERSON_COLORS[2])
add_script_block("""
The login endpoint handles POST /api/auth/login.php.

The process is: we read the Student ID and password from the request body, search the database for that student ID, and if found we use PHP's built-in password_verify() to check the password against the hashed version in the database — we never store raw passwords. If everything matches, we store the user's ID, name, and role in the session.

Sessions are how we remember who you are across requests. PHP creates a session cookie in your browser, and every subsequent request carries that cookie.
""", 2)

add_paragraph('Open: backend/api/auth/register.php', bold=True, size=10, color=PERSON_COLORS[2])
add_script_block("""
The registration endpoint is more complex. Look at these validation checks — we verify the password has at least 8 characters, then use preg_match — a PHP regex function — to check for an uppercase letter and a special character.

These exact same regex patterns are used in the frontend checklist. That consistency is important: the backend is the final authority, but the frontend shows users what is required before they even submit.

We hash the password using password_hash() with PHP's PASSWORD_DEFAULT — currently that is bcrypt. Even if someone stole our database, they would only see scrambled text, never the actual passwords.
""", 2)

add_paragraph('Open: backend/api/reservations/create.php', bold=True, size=10, color=PERSON_COLORS[2])
add_script_block("""
This is probably our most complex file. When you click Confirm Reservation, this runs and validates many things in sequence.

One — the booking time must be in the future. We create a DateTime object and compare it to now.
Two — guests cannot book Shade A tables.
Three — a capacity check: we query how many seats are already booked and make sure there is room for your group.
Four — we check for time conflicts. If another reservation exists at the same venue within 60 minutes of your chosen time, we reject it. This prevents double-booking.
Five — guests are limited to one active reservation at a time.

Only after all these checks pass do we insert the reservation. The security design here is critical: we never trust the frontend. Even if someone bypassed our React interface and sent a direct HTTP request, the backend would still enforce every rule.

I'll now hand over to [Person 4], who will explain how the database is structured.
""", 2)

add_rule()

# ════════════════════════════════════════════════════════════════════════════
#  PERSON 4 — DATABASE
# ════════════════════════════════════════════════════════════════════════════
doc.add_page_break()
person_banner(4, 'Database Admin', 'MySQL Schema & Data Flow', PERSON_COLORS[3])

add_heading('What to Show on Screen', 3, PERSON_COLORS[3])
add_bullet('phpMyAdmin at http://localhost/phpmyadmin')
add_bullet('backend/database/schema.sql open in VS Code')
add_bullet('Click through: users table → venues table → reservations table')

add_heading('Speaking Script', 3, PERSON_COLORS[3])

add_paragraph('Show: phpMyAdmin — university_seat_reservation database', bold=True, size=10, color=PERSON_COLORS[3])
add_script_block("""
Thank you, [Person 3]. My name is [Name] and I will explain our database — where all the data is stored.

We use MySQL as our database, managed through phpMyAdmin, which comes with XAMPP. Our database is called university_seat_reservation and it has exactly three tables.
""", 3)

add_paragraph('Show: users table', bold=True, size=10, color=PERSON_COLORS[3])
add_script_block("""
The first table is users. It stores everyone with an account.

The key columns: student_id — like STRATH/2023/1234 — marked UNIQUE so no two people can have the same ID. We have full_name and email. The password_hash column stores the bcrypt-hashed password, not the real password. The role column is an ENUM — it can only be one of three values: student, staff, or guest. The database enforces this — you literally cannot insert any other value.

We seeded one demo user: Student ID STU001, name Demo Student, with the password password123. In the table you can see that password stored as a long hash string — that is bcrypt in action.
""", 3)

add_paragraph('Show: venues table', bold=True, size=10, color=PERSON_COLORS[3])
add_script_block("""
The second table is venues — the physical tables on campus.

We have 31 venues: 6 in Shade A — labelled A1 to A6 — restricted to students and staff only; and 25 in Shade B — labelled B1 to B25 — open to everyone including guests.

Each venue has a capacity of 6. There is also an is_active flag — if a table is closed for maintenance, setting this to false removes it from the booking map.
""", 3)

add_paragraph('Show: reservations table', bold=True, size=10, color=PERSON_COLORS[3])
add_script_block("""
The third table is reservations — the heart of the system.

The design here is the foreign keys: user_id references the users table, and venue_id references the venues table. This creates relationships. If a user is deleted, all their reservations are automatically deleted too — that is the ON DELETE CASCADE rule.

The status column is another ENUM: active, cancelled, or completed. When a user cancels a reservation we do not delete the row — we just change the status to cancelled. This preserves history.

We also have two indexes — idx_user_status and idx_venue_time. Indexes make lookups faster. Without them, the database would scan every row to find your reservations. With them, it jumps directly to the right rows.

I'll now hand over to [Person 5], who will show how all these parts work together as a complete system.
""", 3)

add_rule()

# ════════════════════════════════════════════════════════════════════════════
#  PERSON 5 — SYSTEM FLOW & CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
doc.add_page_break()
person_banner(5, 'Systems Integrator', 'Full Flow & Conclusion', PERSON_COLORS[4])

add_heading('What to Show on Screen', 3, PERSON_COLORS[4])
add_bullet('Running app — demonstrate a full booking live')
add_bullet('frontend/src/api.js briefly visible')

add_heading('Speaking Script', 3, PERSON_COLORS[4])

add_script_block("""
Thank you, [Person 4]. My name is [Name], and I am going to bring everything together by showing you how the full system works end-to-end — from the moment you open the browser to a confirmed reservation.
""", 4)

add_paragraph('Step 1 — Landing Page', bold=True, size=10, color=PERSON_COLORS[4])
add_script_block("""
The user opens the browser at localhost:5173. The first thing React does is check: is this person already logged in? It sends a GET request to me.php on the backend, which checks the session. If not logged in, it returns authenticated: false and the landing page is shown. If already logged in, the page automatically redirects to the dashboard.
""", 4)

add_paragraph('Step 2 — Login (open frontend/src/api.js)', bold=True, size=10, color=PERSON_COLORS[4])
add_script_block("""
The user clicks Sign In and arrives at the login form. When they submit, the request function in api.js handles it — you can see it adds the right headers, attaches the session cookie with credentials: include, and handles errors consistently across all requests.

The login request goes to localhost:8000/api/auth/login.php. PHP checks the credentials, sets the session, and responds with the user's details. React receives the response and navigates to /home.
""", 4)

add_paragraph('Step 3 — Dashboard Loads', bold=True, size=10, color=PERSON_COLORS[4])
add_script_block("""
The Home page mounts and immediately fires two requests simultaneously — venues.list() and reservations.list() — using Promise.all, so both load in parallel and the page is ready as fast as possible.

The venues response includes real-time availability calculated by the database: how many seats are booked versus available. The SeatMap renders this visually.
""", 4)

add_paragraph('Step 4 — Making a Reservation (live demo)', bold=True, size=10, color=PERSON_COLORS[4])
add_script_block("""
[Click a table on screen] When I click this table it highlights in blue. The ReservationForm on the right recognises the selection.

I choose a group size and a booking time — it must be in the future, the calendar input won't let me pick the past.

I hit Confirm. The frontend sends a POST request with venue_id, group_size, and booking_time to reservations/create.php.

The backend runs through all its checks — future time, capacity, role restrictions, no time conflicts. If everything passes, it inserts the reservation into MySQL.

Back in the browser, refreshData() fetches the updated venue list. Watch the table — the remaining count drops. The reservation also appears in the list on the right.
""", 4)

add_heading('Design Highlights', 3, PERSON_COLORS[4])
add_script_block("""
I want to highlight a few design decisions we are proud of.

Security is layered. Password rules are checked in real time in the frontend, then validated again in PHP before anything is saved, and passwords are stored as bcrypt hashes. Three separate layers.

The backend is the final authority. Even if someone bypassed the entire React interface and sent a raw HTTP request, the PHP code would still enforce every rule.

The API is a clean boundary. The api.js file has about 60 lines of code but defines the entire contract between frontend and backend. All requests go through one shared function with consistent error handling.

The database is normalised. Three tables, proper foreign keys, proper indexes, ENUMs for constrained fields. Simple but correct.
""", 4)

add_heading('Conclusion', 3, PERSON_COLORS[4])
add_script_block("""
To summarise what we built: a student walks in, opens the app, sees every available table in real time, books a spot in 30 seconds, and that reservation is immediately visible to everyone else.

We used React and Vite for a fast, modern frontend; PHP for a secure backend API; and MySQL for reliable data storage. The whole system starts with a single command — npm run dev.

We are proud of the work we put into this and we are happy to answer any questions. Thank you.
""", 4)

# ════════════════════════════════════════════════════════════════════════════
#  TRANSITION LINES PAGE
# ════════════════════════════════════════════════════════════════════════════
doc.add_page_break()
add_heading('Transition Lines — Quick Reference', 2)
transitions = [
    ('Person 1 → 2', '"I\'ll now hand it over to [Name], who will take you through the frontend in detail."'),
    ('Person 2 → 3', '"I\'ll now hand over to [Name], who will explain what happens on the server side."'),
    ('Person 3 → 4', '"I\'ll now hand over to [Name], who\'ll explain how the database is structured."'),
    ('Person 4 → 5', '"I\'ll now hand over to [Name], who will show how all these parts work together."'),
    ('Person 5 → Q&A', '"We are happy to answer any questions. Thank you."'),
]
t = doc.add_table(rows=len(transitions), cols=2)
t.style = 'Table Grid'
for i, (hand, line) in enumerate(transitions):
    t.rows[i].cells[0].text = hand
    t.rows[i].cells[1].text = line
    t.rows[i].cells[0].paragraphs[0].runs[0].bold = True
    t.rows[i].cells[0].paragraphs[0].runs[0].font.color.rgb = NAVY

doc.save('Presentation_Script.docx')
print('[OK]  Presentation_Script.docx saved')


# ── POWERPOINT ───────────────────────────────────────────────────────────────
from pptx import Presentation as PPT
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN

prs = PPT()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

PNAVY  = PRGB(0x1A, 0x3C, 0x6D)
PWHITE = PRGB(0xFF, 0xFF, 0xFF)
PGRAY  = PRGB(0xF0, 0xF4, 0xFF)
PMUTED = PRGB(0x64, 0x74, 0x8B)

P_COLORS = [
    PRGB(0x1A, 0x3C, 0x6D),
    PRGB(0x16, 0x5F, 0x9F),
    PRGB(0x0E, 0x7C, 0x65),
    PRGB(0x7C, 0x2D, 0x12),
    PRGB(0x4C, 0x1D, 0x95),
]

BLANK = prs.slide_layouts[6]   # completely blank

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, size=18, bold=False, color=PRGB(0,0,0),
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def slide_header(slide, person_num, title, subtitle, color):
    # Left accent bar
    add_rect(slide, 0, 0, 0.18, 7.5, fill=color)
    # Top area bg
    add_rect(slide, 0.18, 0, 13.15, 1.5, fill=PNAVY)
    # Person badge
    badge = slide.shapes.add_shape(1, Inches(0.35), Inches(0.25), Inches(1.1), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f'P{person_num}'
    run.font.bold  = True
    run.font.size  = Pt(14)
    run.font.color.rgb = PWHITE
    # Title
    add_text_box(slide, title, 1.6, 0.2, 10.8, 0.75,
                 size=28, bold=True, color=PWHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, subtitle, 1.6, 0.9, 10.8, 0.45,
                 size=14, color=PRGB(0xBF, 0xDB, 0xFE), align=PP_ALIGN.LEFT)

def bullet_slide(slide, items, start_y=1.7, color=PNAVY):
    """items: list of (text, is_sub)"""
    y = start_y
    for text, is_sub in items:
        indent = 0.6 if is_sub else 0.3
        dot_x  = 0.55 if is_sub else 0.35
        dot_size = 0.06 if is_sub else 0.09
        dot = slide.shapes.add_shape(9, Inches(dot_x + 0.18), Inches(y + 0.08),
                                     Inches(dot_size), Inches(dot_size))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color if not is_sub else PMUTED
        dot.line.fill.background()
        txt_size = 16 if not is_sub else 13
        txt_color = PRGB(0x0F, 0x17, 0x2A) if not is_sub else PMUTED
        add_text_box(slide, text, 0.18 + indent, y, 12.2 - indent, 0.38,
                     size=txt_size, color=txt_color)
        y += 0.42 if not is_sub else 0.36
    return y

def code_box(slide, code_lines, y, color):
    h = 0.32 * len(code_lines) + 0.2
    add_rect(slide, 0.35, y, 12.6, h, fill=PRGB(0xF0, 0xF4, 0xFF), line=color)
    txt = '\n'.join(code_lines)
    add_text_box(slide, txt, 0.5, y + 0.05, 12.3, h - 0.1,
                 size=11, color=PRGB(0x1E, 0x40, 0xAF), italic=False)


# ════════ SLIDE 1 — COVER ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PNAVY)
add_rect(sl, 0, 0, 0.45, 7.5, fill=P_COLORS[0])
add_rect(sl, 0.45, 2.5, 12.88, 0.08, fill=PRGB(0x93, 0xC5, 0xFD))
add_text_box(sl, 'UNIVERSITY SEAT\nRESERVATION SYSTEM',
             1.0, 0.9, 11.5, 1.9, size=42, bold=True,
             color=PWHITE, align=PP_ALIGN.LEFT)
add_text_box(sl, 'Group Presentation · Full-Stack Web Application',
             1.0, 2.75, 11.5, 0.55, size=18,
             color=PRGB(0xBF, 0xDB, 0xFE), align=PP_ALIGN.LEFT)

persons = [
    ('P1', 'Project Lead',       P_COLORS[0]),
    ('P2', 'Frontend Developer', P_COLORS[1]),
    ('P3', 'Backend Developer',  P_COLORS[2]),
    ('P4', 'Database Admin',     P_COLORS[3]),
    ('P5', 'Systems Integrator', P_COLORS[4]),
]
for i, (label, role, clr) in enumerate(persons):
    x = 1.0 + i * 2.38
    add_rect(sl, x, 3.55, 2.1, 1.0, fill=PRGB(0x1E, 0x40, 0x6F))
    add_rect(sl, x, 3.55, 2.1, 0.12, fill=clr)
    add_text_box(sl, label, x + 0.08, 3.68, 0.55, 0.38,
                 size=14, bold=True, color=clr)
    add_text_box(sl, role, x + 0.08, 4.05, 1.95, 0.45,
                 size=11, color=PRGB(0xBF, 0xDB, 0xFE))

add_text_box(sl, 'Demo login:  STU001  /  password123     |     npm run dev  to start',
             1.0, 6.7, 11.5, 0.4, size=11,
             color=PRGB(0x60, 0x82, 0xAF), align=PP_ALIGN.LEFT)


# ════════ SLIDE 2 — TECH STACK ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
add_rect(sl, 0, 0, 13.33, 1.1, fill=PNAVY)
add_rect(sl, 0, 0, 0.18, 7.5, fill=P_COLORS[0])
add_text_box(sl, 'Technology Stack & Project Structure',
             0.4, 0.15, 12.5, 0.75, size=26, bold=True,
             color=PWHITE, align=PP_ALIGN.LEFT)

boxes = [
    ('FRONTEND', 'React 19 + Vite 8\nReact Router DOM v7\nTailwind CSS\nLucide React (icons)', P_COLORS[1]),
    ('BACKEND',  'PHP 8.0\nPHP Built-in Server\nPDO for MySQL\nSession-based Auth', P_COLORS[2]),
    ('DATABASE', 'MySQL via XAMPP\nphpMyAdmin\n3 Tables · 31 Venues\nBcrypt passwords', P_COLORS[3]),
    ('TOOLING',  'XAMPP (MySQL)\nnpm + concurrently\nVite dev server\n.env config', P_COLORS[4]),
]
for i, (title, body, clr) in enumerate(boxes):
    x = 0.4 + i * 3.2
    add_rect(sl, x, 1.35, 3.0, 2.8, fill=PWHITE, line=clr)
    add_rect(sl, x, 1.35, 3.0, 0.35, fill=clr)
    add_text_box(sl, title, x + 0.12, 1.38, 2.8, 0.3,
                 size=12, bold=True, color=PWHITE)
    add_text_box(sl, body, x + 0.12, 1.78, 2.8, 2.3,
                 size=12, color=PRGB(0x0F, 0x17, 0x2A))

add_text_box(sl, 'ROUTES', 0.4, 4.35, 12.5, 0.35, size=13, bold=True, color=PNAVY)
routes = [('/ — Landing Page', P_COLORS[0]), ('/login — Login', P_COLORS[1]),
          ('/signup — Register', P_COLORS[2]), ('/home — Dashboard', P_COLORS[3])]
for i, (r, clr) in enumerate(routes):
    x = 0.4 + i * 3.2
    add_rect(sl, x, 4.75, 3.0, 0.55, fill=PRGB(0xEF, 0xF6, 0xFF), line=clr)
    add_text_box(sl, r, x + 0.15, 4.8, 2.8, 0.45, size=13, color=PNAVY)


# ════════ SLIDE 3 — PERSON 1: SYSTEM OVERVIEW ════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 1, 'Introduction & System Overview', 'Project Lead', P_COLORS[0])
bullet_slide(sl, [
    ('Problem: Students waste time searching for available seats on campus', False),
    ('Solution: A web app that shows real-time seat availability and lets you book instantly', False),
    ('Open the browser → see all tables → click one → confirm booking → done', True),
    ('Three user roles, each with different access levels:', False),
    ('Student / Staff — full access to Shade A and B, group up to 6', True),
    ('Guest — Shade B only, group up to 3, one reservation at a time', True),
    ('Demo the live app: localhost:5173', False),
], color=P_COLORS[0])
add_text_box(sl, 'KEY MESSAGE: Simple problem → clean solution → real technology',
             0.35, 6.8, 12.6, 0.45, size=12, bold=True,
             color=PWHITE, align=PP_ALIGN.CENTER)
add_rect(sl, 0.35, 6.75, 12.63, 0.55, fill=P_COLORS[0])
add_text_box(sl, 'KEY MESSAGE: Simple problem → clean solution → real technology',
             0.35, 6.8, 12.6, 0.4, size=12, bold=True,
             color=PWHITE, align=PP_ALIGN.CENTER)


# ════════ SLIDE 4 — PERSON 2: FRONTEND (routing + auth) ══════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 2, 'Frontend — Pages & Components', 'Frontend Developer', P_COLORS[1])
bullet_slide(sl, [
    ('main.jsx — Router entry point, defines 4 routes', False),
    ('Landing.jsx — Homepage: hero image, navbar, learn-more section', False),
    ('Login.jsx — Student ID + password form, live password checklist (informational)', False),
    ('Signup.jsx — Full registration form, checklist ENFORCED, button disabled until valid', False),
    ('Home.jsx — Protected dashboard: auth check → load venues + reservations in parallel', False),
    ('Components/SeatMap.jsx — Visual seat map: white=available, blue=selected, grey=full, amber=restricted', False),
    ('Components/ReservationForm.jsx — Group size, time picker, selected seat badge', False),
    ('Components/ReservationList.jsx — Edit/cancel reservations inline', False),
    ('Components/PasswordChecklist.jsx — Shared RULES constant used by both Login and Signup', False),
], color=P_COLORS[1])


# ════════ SLIDE 5 — PERSON 2: FRONTEND (api + live checklist) ════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 2, 'Frontend — API Layer & Password Validation', 'Frontend Developer', P_COLORS[1])
add_text_box(sl, 'api.js — Single file connecting frontend to backend:', 0.35, 1.7, 12.6, 0.38,
             size=15, bold=True, color=PNAVY)
code_box(sl, [
    "const API_BASE = 'http://localhost:8000/api'",
    "async function request(path, options) { ... credentials: 'include' ... }",
    "export const auth     = { login, register, logout, me }",
    "export const venues   = { list }",
    "export const reservations = { list, create, update, cancel }",
], 2.15, P_COLORS[1])
add_text_box(sl, 'Live Password Checklist — RULES array (same regex as PHP backend):', 0.35, 4.25, 12.6, 0.38,
             size=15, bold=True, color=PNAVY)
code_box(sl, [
    "{ key: 'length',    test: (pw) => pw.length >= 8          }",
    "{ key: 'uppercase', test: (pw) => /[A-Z]/.test(pw)        }",
    "{ key: 'special',   test: (pw) => /[^A-Za-z0-9]/.test(pw) }",
], 4.7, P_COLORS[1])
add_text_box(sl, 'Login: checklist shown, NOT enforced  ·  Signup: checklist enforced, button disabled',
             0.35, 6.35, 12.6, 0.4, size=13, italic=True, color=PMUTED)


# ════════ SLIDE 6 — PERSON 3: BACKEND OVERVIEW ═══════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 3, 'Backend — PHP APIs & Server Logic', 'Backend Developer', P_COLORS[2])
bullet_slide(sl, [
    ('config/database.php — Loaded by EVERY PHP file; provides:', False),
    ('Env variable loading (.env) — DB creds never in code', True),
    ('CORS headers — allows frontend on :5173 to call backend on :8000', True),
    ('4 helpers: getDb()  jsonResponse()  errorResponse()  requireAuth()', True),
    ('api/auth/login.php — Verifies student_id + password_verify() → sets session', False),
    ('api/auth/register.php — Validates input + regex + duplicate check → bcrypt hash → insert', False),
    ('api/auth/me.php — Checks session; returns {authenticated: true/false}', False),
    ('api/auth/logout.php — Clears session', False),
    ('api/venues/index.php — Returns venues with real-time remaining capacity', False),
    ('api/reservations/ — create · update · cancel · index', False),
], color=P_COLORS[2])


# ════════ SLIDE 7 — PERSON 3: RESERVATION LOGIC ══════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 3, 'Backend — Reservation Business Rules', 'Backend Developer', P_COLORS[2])
add_text_box(sl, 'reservations/create.php — Checks in sequence:', 0.35, 1.65, 12.6, 0.38,
             size=15, bold=True, color=PNAVY)
checks = [
    '1.  booking_time > now  (DateTime comparison)',
    '2.  Venue exists and is_active = true',
    '3.  Guest + Shade A → reject (403)',
    '4.  group_size ≤ maxGroupSize  (6 for students/staff, 3 for guests)',
    '5.  group_size ≤ remaining capacity  (capacity - SUM of active reservations)',
    '6.  Guest with existing active reservation → reject (one at a time)',
    '7.  No other reservation within 60 minutes at same venue → reject (409)',
    '✓   All checks pass → INSERT into reservations table',
]
y = 2.1
for chk in checks:
    clr = P_COLORS[2] if chk.startswith('✓') else PRGB(0x0F, 0x17, 0x2A)
    add_text_box(sl, chk, 0.55, y, 12.3, 0.38, size=13, color=clr,
                 bold=chk.startswith('✓'))
    y += 0.42
add_rect(sl, 0.35, 6.7, 12.63, 0.5, fill=P_COLORS[2])
add_text_box(sl, 'Security: backend enforces ALL rules — frontend is just a convenience layer',
             0.5, 6.75, 12.4, 0.38, size=12, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)


# ════════ SLIDE 8 — PERSON 4: DATABASE ═══════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 4, 'Database — MySQL Schema & Design', 'Database Admin', P_COLORS[3])

tables_info = [
    ('users', ['id  ·  student_id (UNIQUE)', 'full_name  ·  email (UNIQUE)',
               'password_hash  (bcrypt)', 'role ENUM(student|staff|guest)',
               'created_at  ·  updated_at'], P_COLORS[3]),
    ('venues', ['id  ·  table_id (UNIQUE)', 'shade ENUM(A|B)',
                'capacity (default 6)', 'is_active BOOLEAN',
                'Seeded: 6×A + 25×B = 31 total'], P_COLORS[2]),
    ('reservations', ['id  ·  user_id (FK→users)', 'venue_id (FK→venues)',
                      'group_size  ·  booking_time',
                      'status ENUM(active|cancelled|completed)',
                      'ON DELETE CASCADE  ·  2 indexes'], P_COLORS[1]),
]
for i, (name, cols, clr) in enumerate(tables_info):
    x = 0.35 + i * 4.28
    add_rect(sl, x, 1.55, 4.0, 3.5, fill=PWHITE, line=clr)
    add_rect(sl, x, 1.55, 4.0, 0.38, fill=clr)
    add_text_box(sl, name, x + 0.12, 1.58, 3.8, 0.3, size=14, bold=True, color=PWHITE)
    col_txt = '\n'.join(cols)
    add_text_box(sl, col_txt, x + 0.12, 2.02, 3.8, 3.0, size=11.5,
                 color=PRGB(0x0F, 0x17, 0x2A))

add_text_box(sl, 'Key Design Choices', 0.35, 5.2, 12.6, 0.38, size=14, bold=True, color=PNAVY)
bullet_slide(sl, [
    ('Soft deletes: cancelled reservations keep status=cancelled, history preserved', False),
    ('ENUMs enforce valid values at DB level (not just application level)', False),
    ('Indexes on (user_id, status) and (venue_id, booking_time) for fast queries', False),
], start_y=5.65, color=P_COLORS[3])


# ════════ SLIDE 9 — PERSON 5: SYSTEM FLOW ════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 5, 'Complete System Flow — End to End', 'Systems Integrator', P_COLORS[4])

steps = [
    ('1  Browser opens localhost:5173',   'React loads → auth.me() → {authenticated:false} → show Landing',    P_COLORS[0]),
    ('2  User clicks Sign In / Submit',   'POST /auth/login.php → PHP verifies hash → sets session → /home',    P_COLORS[1]),
    ('3  Dashboard mounts',              'Promise.all([venues.list(), reservations.list()]) → SeatMap renders', P_COLORS[2]),
    ('4  User clicks a table',           'onSeatClick → selectedSeats state → ReservationForm activates',       P_COLORS[3]),
    ('5  User confirms reservation',     'POST /reservations/create.php → 7 checks → INSERT → refreshData()',   P_COLORS[4]),
    ('6  Table updates live',            'remaining count drops on screen · reservation appears in list',        P_COLORS[0]),
]
y = 1.7
for label, desc, clr in steps:
    add_rect(sl, 0.35, y, 4.5, 0.5, fill=clr)
    add_text_box(sl, label, 0.48, y + 0.08, 4.35, 0.35, size=12, bold=True, color=PWHITE)
    add_rect(sl, 4.85, y, 8.1, 0.5, fill=PWHITE, line=clr)
    add_text_box(sl, desc, 4.98, y + 0.08, 7.9, 0.38, size=11.5, color=PRGB(0x0F, 0x17, 0x2A))
    y += 0.65


# ════════ SLIDE 10 — PERSON 5: DESIGN HIGHLIGHTS ═════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=PRGB(0xF8, 0xFA, 0xFC))
slide_header(sl, 5, 'Design Highlights & Conclusion', 'Systems Integrator', P_COLORS[4])
highlights = [
    ('Security is layered',         'Frontend checklist → PHP validation → bcrypt storage. Three independent layers.',  P_COLORS[0]),
    ('Backend is final authority',  'All 7 reservation rules enforced in PHP. Bypassing React changes nothing.',         P_COLORS[1]),
    ('Clean API boundary',          'api.js: ~60 lines, one shared request() function, consistent error handling.',      P_COLORS[2]),
    ('Normalised database',         '3 tables, foreign keys, ENUMs, indexes — simple, correct, scalable.',              P_COLORS[3]),
    ('Single start command',        'npm run dev — starts PHP on :8000 and Vite on :5173 concurrently.',                P_COLORS[4]),
]
y = 1.7
for title, desc, clr in highlights:
    add_rect(sl, 0.35, y, 0.12, 0.5, fill=clr)
    add_text_box(sl, title, 0.58, y + 0.03, 3.5, 0.25, size=13, bold=True, color=clr)
    add_text_box(sl, desc, 0.58, y + 0.26, 12.1, 0.25, size=12, color=PRGB(0x0F, 0x17, 0x2A))
    y += 0.65

add_rect(sl, 0.35, 6.65, 12.63, 0.6, fill=PNAVY)
add_text_box(sl, 'Thank you — we are happy to answer any questions',
             0.5, 6.68, 12.4, 0.48, size=15, bold=True,
             color=PWHITE, align=PP_ALIGN.CENTER)

prs.save('Presentation_Slides.pptx')
print('[OK]  Presentation_Slides.pptx saved')
print()
print('Both files are ready in the Presentation/ folder.')
